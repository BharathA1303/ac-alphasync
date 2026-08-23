"""
Extract plain text from uploaded lesson material files (PDF/DOCX/PPTX/MD)
so it can be used as AI context when generating MCQ assessment questions.
"""

import logging

logger = logging.getLogger(__name__)

MAX_EXTRACTED_CHARS = 20000


def extract_text(file_path: str, file_type: str) -> str:
    """Best-effort text extraction. Returns '' on any failure — never raises."""
    try:
        if file_type == "pdf":
            return _extract_pdf(file_path)
        if file_type == "docx":
            return _extract_docx(file_path)
        if file_type == "pptx":
            return _extract_pptx(file_path)
        if file_type == "md":
            return _extract_text_file(file_path)
    except Exception:
        logger.exception("Failed to extract text from %s (%s)", file_path, file_type)
    return ""


def _extract_pdf(file_path: str) -> str:
    from pypdf import PdfReader

    reader = PdfReader(file_path)
    chunks = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(chunks)[:MAX_EXTRACTED_CHARS]


def _extract_docx(file_path: str) -> str:
    from docx import Document

    doc = Document(file_path)
    chunks = [p.text for p in doc.paragraphs if p.text]
    return "\n".join(chunks)[:MAX_EXTRACTED_CHARS]


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks)[:MAX_EXTRACTED_CHARS]


def _extract_text_file(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()[:MAX_EXTRACTED_CHARS]
